#!/usr/bin/env python3
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches

# let's automate powerpoint

name = input("presentation name? : ")

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = name #"Hello, World!"
subtitle.text = "python-pptx was here!"

#---

title_slide_layout1 = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(title_slide_layout1)

titleFABIsSlide = slide2.shapes.title
subtitleFABIsSlide = slide2.placeholders[0]

titleFABIsSlide.text = "Hi!!!"
subtitleFABIsSlide.text = "FABI was here!"

left = Inches(1)
top = Inches(3)
pictureFile = './user-avatar.jpeg'
slide2.shapes.add_picture(pictureFile, left, top)


#---

title_slide_layout1 = prs.slide_layouts[5]
slide3 = prs.slides.add_slide(title_slide_layout1)

slideShape3 = slide3.shapes

slideShape3.title.text = 'Adding a table'

rows2 = 4
cols2 = 2
left2 = top2 = Inches(2.0)
width2 = Inches(6.0)
height2 = Inches(0.8)

table = slideShape3.add_table(rows2, cols2, left2, top2, width2, height2).table

# set column widths
table.columns[0].width2 = Inches(2.0)
table.columns[1].width2 = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'

cell = table.cell(1, 0)

table.cell(2, 0).text = 'FABI'
table.cell(2, 1).text = str(cell.text) + str(table.cell(2, 0).text)

table.cell(3, 0).text = '2'
table.cell(3, 1).text = str(int(table.cell(3, 0).text) + 1)

prs.save('test.pptx')
